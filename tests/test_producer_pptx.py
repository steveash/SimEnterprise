"""Tests for the PowerPoint (``pptx``) producer (ARCHITECTURE.md §4, §16, esim-8db34f7e).

Acceptance: a kickoff deck ``.pptx`` with real slides opens/validates, and the
presentation kind can be *rebound* from the markdown default to ``pptx``. Also
covers templated grounding (title/team/references slides), mentions, reference
edges, the single repair pass + validation issue, and determinism.

The "opens/validates" check is done in-process by re-opening the produced bytes
with python-pptx (the same library PowerPoint-compatible tools use to parse a
deck); a malformed package raises on open.
"""

from __future__ import annotations

import io
from collections.abc import Mapping, Sequence
from datetime import UTC, datetime
from typing import Any

import pytest
from enterprise_sim.core.events import Deliverable, Event
from enterprise_sim.core.llm import LLMClient
from enterprise_sim.core.llm.backends import estimate_tokens
from enterprise_sim.core.llm.prompt import Prompt
from enterprise_sim.core.llm.types import Completion, TokenUsage
from enterprise_sim.core.registry import PRODUCERS, BindingMap
from enterprise_sim.core.world import Node, World
from enterprise_sim.producers import (
    PptxProducer,
    apply_to_world,
    build_kickoff_deck,
    mention_records,
)
from enterprise_sim.producers.pptx import Slide, sample_deck
from pptx import Presentation

_T0 = datetime(2026, 6, 1, 9, 0, tzinfo=UTC)
_T_EVENT = datetime(2026, 6, 12, 14, 0, tzinfo=UTC)


def _person(node_id: str, name: str, *, aliases: list[str] | None = None) -> Node:
    return Node(node_id, "Person", _T0, props={"name": name}, aliases=aliases or [])


def _artifact(node_id: str, title: str) -> Node:
    return Node(node_id, "Artifact", _T0, props={"title": title, "kind": "design_doc"})


def _world() -> World:
    world = World()
    world.add_node(_person("person:ada", "Ada Lovelace", aliases=["Ada"]))
    world.add_node(_person("person:alan", "Alan Turing"))
    world.add_node(Node("project:payments", "Project", _T0, props={"name": "Payments Platform"}))
    world.add_node(_artifact("artifact:design-pay", "Payments Design"))
    return world


def _event() -> Event:
    return Event(
        id="evt:kickoff-w12",
        type="DeliverableDrafted",
        timestamp=_T_EVENT,
        actors={"author": ["person:ada"], "reviewers": ["person:alan"]},
        initiative="init:payments",
        project="project:payments",
        subjects=["project:payments"],
        deliverable=Deliverable(kind="kickoff_deck", medium="presentation"),
        payload={"topic": "payments rollout", "title": "Payments Platform Kickoff"},
    )


def _fake_client() -> LLMClient:
    from enterprise_sim.core.llm import LLMConfig

    return LLMClient.from_config(LLMConfig(backend="fake", cache_enabled=False))


def _produce(client: LLMClient | None = None) -> Any:
    view = _world().projection(at=_T_EVENT)
    return PptxProducer().produce(_event(), view, client or _fake_client())


def _reopen(produced: Any) -> Any:
    """Re-open the produced ``.pptx`` bytes — raises if the package is malformed."""
    assert produced.binary_body is not None
    return Presentation(io.BytesIO(produced.binary_body))


def _slide_titles(prs: Any) -> list[str]:
    return [s.shapes.title.text for s in prs.slides if s.shapes.title is not None]


# -- the deck opens / validates (the acceptance criterion) ------------------


def test_produce_renders_pptx_that_opens() -> None:
    produced = _produce()
    assert produced.fmt == "pptx"
    assert produced.path == "artifacts/artifact-kickoff-w12.pptx"
    assert produced.is_binary and produced.binary_body
    # Re-opening the bytes proves the package is a valid, parseable .pptx.
    prs = _reopen(produced)
    assert len(prs.slides) >= 4  # real slides, not a stub


def test_deck_has_expected_slide_sequence() -> None:
    prs = _reopen(_produce())
    titles = _slide_titles(prs)
    assert titles[0] == "Payments Platform Kickoff"  # templated title slide
    assert "Agenda" in titles
    assert "Overview" in titles
    assert "Team & Roles" in titles
    assert "References" in titles  # candidate exists → citations slide present


def test_title_slide_is_templated_and_grounded() -> None:
    prs = _reopen(_produce())
    title_slide = prs.slides[0]
    subtitle = title_slide.placeholders[1].text
    # Date + author are templated from bound roles, never generated (§16.2.2).
    assert "2026-06-12" in subtitle
    assert "Ada Lovelace" in subtitle


def test_team_slide_lists_roles_from_bound_roles() -> None:
    prs = _reopen(_produce())
    team = next(s for s in prs.slides if s.shapes.title and s.shapes.title.text == "Team & Roles")
    text = "\n".join(p.text for shape in team.shapes for p in _paras(shape))
    assert "Ada Lovelace — Lead" in text
    assert "Alan Turing — Reviewer" in text


def _paras(shape: Any) -> list[Any]:
    return list(shape.text_frame.paragraphs) if shape.has_text_frame else []


# -- KG facts: mentions, edges, metadata ------------------------------------


def test_produce_records_mentions() -> None:
    produced = _produce()
    entities = {m.entity_id for m in produced.mentions}
    assert "person:ada" in entities
    assert "person:alan" in entities
    # Every mention's locator addresses the outline body exactly.
    for mention in produced.mentions:
        span = produced.body[
            mention.locator.offset : mention.locator.offset + mention.locator.length
        ]
        assert span == mention.surface_form


def test_produce_creates_reference_edges_and_slide() -> None:
    produced = _produce()
    ref_edges = [e for e in produced.edges if e.type == "references"]
    assert ref_edges, "fake backend cites a deterministic subset of candidates"
    for edge in ref_edges:
        assert edge.src == produced.artifact_id
        assert edge.dst == "artifact:design-pay"
    assert produced.metadata["references"] == [e.dst for e in ref_edges]
    # The citation is templated into the References slide too.
    prs = _reopen(produced)
    refs = next(s for s in prs.slides if s.shapes.title and s.shapes.title.text == "References")
    text = "\n".join(p.text for shape in refs.shapes for p in _paras(shape))
    assert "artifact:design-pay" in text


def test_produce_creates_authored_reviewed_and_provenance_edges() -> None:
    produced = _produce()
    by_type = {e.type: e for e in produced.edges}
    assert by_type["authored"].src == "person:ada"
    assert by_type["authored"].dst == produced.artifact_id
    assert by_type["reviewed"].src == "person:alan"
    expresses = [e for e in produced.edges if e.type == "expresses"]
    assert any(e.dst == "project:payments" for e in expresses)


def test_artifact_node_records_format_and_slide_count() -> None:
    produced = _produce()
    assert produced.node.type == "Artifact"
    assert produced.node.props["format"] == "pptx"
    assert produced.node.props["slide_count"] == len(_reopen(produced).slides)
    assert produced.node.props["path"] == produced.path


def test_produce_is_deterministic() -> None:
    a = _produce()
    b = _produce()
    assert a.binary_body == b.binary_body  # python-pptx writes a fixed package stamp
    assert a.body == b.body
    assert [e.id for e in a.edges] == [e.id for e in b.edges]


def test_produce_uses_existing_abstract_artifact_subject() -> None:
    event = _event()
    event.subjects = ["artifact:kickoff-deck", "project:payments"]
    produced = PptxProducer().produce(event, _world().projection(at=_T_EVENT), _fake_client())
    assert produced.artifact_id == "artifact:kickoff-deck"


def test_no_candidates_means_no_references_slide() -> None:
    world = World()
    world.add_node(_person("person:ada", "Ada Lovelace"))
    world.add_node(_person("person:alan", "Alan Turing"))
    world.add_node(Node("project:payments", "Project", _T0, props={"name": "Payments Platform"}))
    produced = PptxProducer().produce(_event(), world.projection(at=_T_EVENT), _fake_client())
    assert [e for e in produced.edges if e.type == "references"] == []
    assert "References" not in _slide_titles(_reopen(produced))


def test_apply_to_world_adds_node_and_edges() -> None:
    world = _world()
    produced = PptxProducer().produce(_event(), world.projection(at=_T_EVENT), _fake_client())
    apply_to_world(world, [produced])
    assert world.get_node(produced.artifact_id) is not None
    for edge in produced.edges:
        assert world.get_edge(edge.id) is not None
    assert mention_records([produced]) == [m.to_dict() for m in produced.mentions]


# -- detect + single repair (§16.2.3 / D17) ---------------------------------


class _ScriptedBackend:
    """A backend whose ``generate_content`` returns queued overview bodies in order."""

    name = "scripted"

    def __init__(self, bodies: Sequence[str]) -> None:
        self._bodies = list(bodies)
        self.calls = 0
        self.prompts: list[Prompt] = []

    def generate_structured(
        self, prompt: Prompt, *, schema: Mapping[str, Any], model: str, temperature: float
    ) -> Completion:  # pragma: no cover - unused
        raise NotImplementedError

    def generate_content(
        self,
        prompt: Prompt,
        *,
        candidate_references: Sequence[str],
        model: str,
        temperature: float,
    ) -> Completion:
        self.prompts.append(prompt)
        body = self._bodies[min(self.calls, len(self._bodies) - 1)]
        self.calls += 1
        usage = TokenUsage(output_tokens=estimate_tokens(body))
        return Completion(text=body, usage=usage, model=model)


def _client_with(backend: _ScriptedBackend) -> LLMClient:
    from enterprise_sim.core.llm import LLMConfig

    return LLMClient(backend, config=LLMConfig(backend="scripted", cache_enabled=False))


def test_repair_pass_fixes_unresolved_name() -> None:
    backend = _ScriptedBackend(
        [
            "Ada Lovelace met Grace Hopper to plan the rollout.",  # hallucinated name
            "Ada Lovelace planned the rollout with the team.",  # clean rewrite
        ]
    )
    produced = _produce(_client_with(backend))
    assert backend.calls == 2
    assert produced.issues == []
    assert "Grace Hopper" not in produced.body


def test_unrepaired_name_becomes_validation_issue() -> None:
    backend = _ScriptedBackend(["Ada Lovelace and Grace Hopper shipped it."])
    produced = _produce(_client_with(backend))
    assert backend.calls == 2  # draft + single repair, then give up
    assert len(produced.issues) == 1
    assert produced.issues[0].kind == "unresolved_mention"
    assert "Grace Hopper" in produced.issues[0].details["names"]
    # The deck is kept despite the issue (§16.2.3 / D17) and still opens.
    assert len(_reopen(produced).slides) >= 4


def test_clean_draft_skips_repair() -> None:
    backend = _ScriptedBackend(["Ada Lovelace led the payments kickoff."])
    _produce(_client_with(backend))
    assert backend.calls == 1


# -- build_kickoff_deck primitives ------------------------------------------


def test_build_kickoff_deck_opens_and_is_deterministic() -> None:
    a = build_kickoff_deck(sample_deck())
    b = build_kickoff_deck(sample_deck())
    assert a == b
    prs = Presentation(io.BytesIO(a))
    assert len(prs.slides) == len(sample_deck())


def test_build_kickoff_deck_rejects_empty() -> None:
    with pytest.raises(ValueError, match="at least one slide"):
        build_kickoff_deck([])


def test_content_slide_keeps_its_bullets() -> None:
    data = build_kickoff_deck(
        [Slide(title="Intro"), Slide(title="Points", bullets=("one", "two", "three"))]
    )
    prs = Presentation(io.BytesIO(data))
    points = prs.slides[1]
    body = next(p.text for shape in points.shapes for p in _paras(shape) if p.text == "two")
    assert body == "two"


# -- rebind: presentation kind → pptx producer (§4, D4/D5) ------------------


def test_producer_handles_presentation_kinds() -> None:
    handles = set(PptxProducer().handles)
    assert "kickoff_deck" in handles
    assert "presentation" in handles


def test_pptx_producer_self_registers() -> None:
    # Importing the producer package registered the plugin into the catalog.
    assert "pptx" in PRODUCERS
    assert PRODUCERS.get("pptx").name == "pptx"


def test_binding_rebinds_kickoff_deck_to_pptx() -> None:
    bindings = BindingMap(default="markdown")
    # Default routes presentation kinds to markdown until rebound.
    assert bindings.producer_names("kickoff_deck") == ["markdown"]
    bindings.bind("kickoff_deck", "pptx")
    assert bindings.producer_names("kickoff_deck") == ["pptx"]
    assert bindings.producer_names("status_report") == ["markdown"]  # other kinds untouched
    # Resolving against the live registry yields the real producer plugin.
    resolved = bindings.resolve("kickoff_deck", PRODUCERS)
    assert [p.name for p in resolved] == ["pptx"]

"""The ``pptx`` producer — events → grounded PowerPoint kickoff decks (ARCHITECTURE.md §4, §16).

The v1 :class:`~enterprise_sim.producers.markdown.MarkdownProducer` is the catch-all
the binding map routes every deliverable kind to; this is the first *rebind* (§4,
D4/D5): bind a presentation kind (``kickoff_deck``) to ``pptx`` and the same
format-agnostic event renders a real ``.pptx`` deck instead of markdown, with the
simulator untouched.

Like the markdown producer it is a pure function of ``(event, view)`` and applies
the three grounding layers (D30):

1. **Constrained input** — a :class:`~enterprise_sim.producers.grounding.Roster`
   built from the view restricts which people/entities the deck may name (§16.2.1).
2. **Templated references** — the title slide's prepared-by line, the date, the
   team-and-roles slide, and the citations slide are filled from the event's bound
   roles + verified ``references_used``, never generated (§16.2.2). Only the
   overview narrative is LLM-written.
3. **Detect + one repair** — the overview prose is run through the mention tagger;
   an unresolved name-like token triggers a *single* repair re-prompt, and anything
   still unresolved becomes a ``validation/issues.jsonl`` entry (§16.2.3 / D17).

The deck itself is built with **python-pptx** (real OOXML slides, not a markdown
stand-in). The :class:`~enterprise_sim.producers.artifact.ProducedArtifact` carries
the ``.pptx`` bytes in ``binary_body`` (what the runner writes) and a plain-text
*outline* in ``body`` — the deterministic textual rendering the mention tagger and
grounding layers reason over. python-pptx writes a fixed package timestamp, so the
bytes are deterministic given a deterministic backend.
"""

from __future__ import annotations

import io
import re
from collections.abc import Sequence
from dataclasses import dataclass
from datetime import datetime

from pptx import Presentation
from pptx.presentation import Presentation as PresentationObj

from enterprise_sim.core.events import Event
from enterprise_sim.core.llm import LLMClient, Prompt, assemble_prompt
from enterprise_sim.core.registry import PRODUCERS
from enterprise_sim.core.world import Edge, Node, WorldView
from enterprise_sim.producers.artifact import (
    Mention,
    ProducedArtifact,
    ValidationIssue,
)
from enterprise_sim.producers.grounding import (
    Roster,
    detect_unresolved_names,
    tag_mentions,
)
from enterprise_sim.producers.markdown import ProducerContext

__all__ = ["PptxProducer", "Slide", "build_kickoff_deck", "sample_deck"]

# KG node/edge type constants this producer reads and writes (§3).
N_ARTIFACT = "Artifact"
E_AUTHORED = "authored"
E_REVIEWED = "reviewed"
E_EXPRESSES = "expresses"
E_REFERENCES = "references"

# Role names in ``event.actors`` mapped to the prepared-by / reviewer templated lines.
_AUTHOR_ROLES = ("author", "authors", "lead", "owner")
_REVIEWER_ROLES = ("reviewer", "reviewers")

# python-pptx default-template layout indices.
_LAYOUT_TITLE = 0  # "Title Slide": centre title + subtitle.
_LAYOUT_CONTENT = 1  # "Title and Content": title + bulleted body placeholder.

# Deliverable kinds this producer renders (presentation decks). ``kickoff_deck`` is
# the kind the binding map rebinds away from the markdown default.
_HANDLES: tuple[str, ...] = ("kickoff_deck", "kickoff", "presentation", "deck")


@dataclass(frozen=True, slots=True)
class Slide:
    """One rendered slide: a ``title`` plus zero or more ``bullets`` (body lines).

    A title-only/section slide has empty ``bullets``. ``subtitle`` is used for the
    opening title slide's centred subtitle; content slides leave it empty.
    """

    title: str
    bullets: tuple[str, ...] = ()
    subtitle: str = ""


class PptxProducer:
    """Render an :class:`Event`'s presentation deliverable to a grounded ``.pptx`` deck.

    Satisfies the :class:`~enterprise_sim.core.registry.plugins.Producer` protocol
    (``name`` / ``formats`` / ``handles``). It is *not* a catch-all — it handles the
    presentation kinds in :data:`_HANDLES`; the binding map rebinds those to it.
    """

    name = "pptx"
    formats: Sequence[str] = ("pptx",)
    handles: Sequence[str] = _HANDLES

    def produce(
        self,
        event: Event,
        view: WorldView,
        client: LLMClient,
        ctx: ProducerContext | None = None,
    ) -> ProducedArtifact:
        """Render ``event`` against ``view`` into a kickoff deck :class:`ProducedArtifact`.

        Returns the ``.pptx`` bytes (``binary_body``), the plain-text outline
        (``body``), the ``Artifact`` node, the ``authored`` / ``reviewed`` /
        ``expresses`` / ``references`` edges, tagged mentions, validation issues, and
        a JSON metadata twin. Nothing is written to disk and ``view`` is not mutated.
        """
        ctx = ctx or ProducerContext()
        roster = Roster.from_worldview(view)

        kind = event.deliverable.kind if event.deliverable else "kickoff_deck"
        medium = event.deliverable.medium if event.deliverable else "presentation"
        title = _title_for(event, kind)
        artifact_id = _artifact_id_for(event)
        path = f"{ctx.artifacts_dir}/{_slug(artifact_id)}.pptx"

        authors = _resolve_people(event, view, _AUTHOR_ROLES)
        reviewers = _resolve_people(event, view, _REVIEWER_ROLES)
        candidates = _candidate_references(view, exclude=artifact_id)
        others = _other_people(view, exclude={n.id for n in authors + reviewers})

        # --- grounded overview prose, then detect + (at most one) repair --------
        overview, references_used, issues = self._generate_grounded(
            client=client,
            event=event,
            kind=kind,
            title=title,
            roster=roster,
            authors=authors,
            reviewers=reviewers,
            candidates=candidates,
            ctx=ctx,
            path=path,
        )

        slides = _build_slides(
            title=title,
            kind=kind,
            event=event,
            authors=authors,
            reviewers=reviewers,
            others=others,
            overview=overview,
            references=references_used,
            view=view,
        )
        binary_body = build_kickoff_deck(slides)
        body = _render_outline(slides)
        mentions = tag_mentions(body, roster, artifact_path=path)

        node = _artifact_node(
            artifact_id=artifact_id,
            kind=kind,
            medium=medium,
            title=title,
            event=event,
            authors=authors,
            reviewers=reviewers,
            path=path,
            slide_count=len(slides),
        )
        edges = _build_edges(
            artifact_id=artifact_id,
            event=event,
            view=view,
            authors=authors,
            reviewers=reviewers,
            references=references_used,
        )
        metadata = _metadata(
            artifact_id=artifact_id,
            kind=kind,
            medium=medium,
            title=title,
            path=path,
            event=event,
            authors=authors,
            reviewers=reviewers,
            references=references_used,
            mentions=mentions,
            issues=issues,
            slide_count=len(slides),
        )
        return ProducedArtifact(
            artifact_id=artifact_id,
            path=path,
            fmt="pptx",
            body=body,
            node=node,
            edges=edges,
            mentions=mentions,
            issues=issues,
            metadata=metadata,
            binary_body=binary_body,
        )

    # -- generation + grounding repair ------------------------------------

    def _generate_grounded(
        self,
        *,
        client: LLMClient,
        event: Event,
        kind: str,
        title: str,
        roster: Roster,
        authors: list[Node],
        reviewers: list[Node],
        candidates: list[str],
        ctx: ProducerContext,
        path: str,
    ) -> tuple[str, tuple[str, ...], list[ValidationIssue]]:
        """Generate the overview narrative, then the detect + single-repair loop (D30.3).

        Returns ``(overview, verified_references, issues)``; mirrors the markdown
        producer's loop so both producers ground identically.
        """
        system = _system_prompt(kind)
        stable = _stable_context(ctx)
        brief = _task_brief(
            event=event,
            kind=kind,
            title=title,
            roster=roster,
            authors=authors,
            reviewers=reviewers,
            candidates=candidates,
        )
        prompt = assemble_prompt(system=system, stable_context=stable, brief=brief)
        result = client.generate_content(prompt, candidate_references=candidates)
        overview = result.content
        references = result.references_used

        unresolved = detect_unresolved_names(overview, roster)
        if unresolved:
            repair_prompt = _repair_prompt(system, stable, brief, unresolved)
            repaired = client.generate_content(repair_prompt, candidate_references=candidates)
            overview = repaired.content
            references = repaired.references_used
            unresolved = detect_unresolved_names(overview, roster)

        issues: list[ValidationIssue] = []
        if unresolved:
            issues.append(
                ValidationIssue(
                    kind="unresolved_mention",
                    message=(
                        "names not resolvable to an in-scope entity survived one "
                        f"repair pass: {', '.join(unresolved)}"
                    ),
                    where=path,
                    details={"names": list(unresolved)},
                )
            )
        return overview, references, issues


# -- deck assembly (python-pptx) --------------------------------------------


def build_kickoff_deck(slides: Sequence[Slide]) -> bytes:
    """Build a real ``.pptx`` deck (as bytes) from ``slides`` using python-pptx.

    The first slide is laid out as the title slide (centre title + subtitle); every
    later slide is a "Title and Content" slide whose bullets fill the body
    placeholder. Output is deterministic for identical input (python-pptx writes a
    fixed package timestamp).
    """
    if not slides:
        raise ValueError("a deck needs at least one slide")
    prs: PresentationObj = Presentation()
    for index, slide in enumerate(slides):
        if index == 0:
            _add_title_slide(prs, slide)
        else:
            _add_content_slide(prs, slide)
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def _add_title_slide(prs: PresentationObj, slide: Slide) -> None:
    layout = prs.slide_layouts[_LAYOUT_TITLE]
    obj = prs.slides.add_slide(layout)
    if obj.shapes.title is not None:
        obj.shapes.title.text = slide.title
    if slide.subtitle:
        obj.placeholders[1].text = slide.subtitle


def _add_content_slide(prs: PresentationObj, slide: Slide) -> None:
    layout = prs.slide_layouts[_LAYOUT_CONTENT]
    obj = prs.slides.add_slide(layout)
    if obj.shapes.title is not None:
        obj.shapes.title.text = slide.title
    if not slide.bullets:
        return
    frame = obj.placeholders[1].text_frame
    frame.text = slide.bullets[0]
    for bullet in slide.bullets[1:]:
        frame.add_paragraph().text = bullet


def _build_slides(
    *,
    title: str,
    kind: str,
    event: Event,
    authors: list[Node],
    reviewers: list[Node],
    others: list[Node],
    overview: str,
    references: Sequence[str],
    view: WorldView,
) -> list[Slide]:
    """Assemble the kickoff deck's slide sequence (templated frame + LLM overview).

    Slides: title (templated) → agenda (fixed) → overview (the only LLM-written
    slide) → team & roles (from bound roles, grounded by construction) → references
    (verified citations, only when present). Every slide but the overview is
    grounded by construction.
    """
    date = event.timestamp.date().isoformat()
    slides: list[Slide] = [
        Slide(title=title, subtitle=_subtitle(kind, date, authors)),
        Slide(title="Agenda", bullets=_AGENDA),
        Slide(title="Overview", bullets=_overview_bullets(overview)),
    ]
    team = _team_bullets(authors, reviewers, others)
    if team:
        slides.append(Slide(title="Team & Roles", bullets=team))
    ref_bullets = _reference_bullets(references, view)
    if ref_bullets:
        slides.append(Slide(title="References", bullets=ref_bullets))
    return slides


_AGENDA: tuple[str, ...] = (
    "Overview",
    "Team & Roles",
    "Next Steps",
)


def _subtitle(kind: str, date: str, authors: list[Node]) -> str:
    human = kind.replace("_", " ").title()
    line = f"{human} · {date}"
    if authors:
        line += f" · Prepared by {_names(authors)}"
    return line


def _overview_bullets(overview: str) -> tuple[str, ...]:
    """Split the overview narrative into slide body lines (one paragraph per bullet)."""
    parts = [seg.strip() for seg in re.split(r"\n+", overview.strip()) if seg.strip()]
    return tuple(parts) if parts else ("(no overview provided)",)


def _team_bullets(
    authors: list[Node], reviewers: list[Node], others: list[Node]
) -> tuple[str, ...]:
    """Templated "Name — Role" lines from the event's bound roles + in-scope people."""
    bullets: list[str] = []
    for node in authors:
        bullets.append(f"{_display_name(node)} — Lead")
    for node in reviewers:
        bullets.append(f"{_display_name(node)} — Reviewer")
    for node in others:
        bullets.append(f"{_display_name(node)} — Contributor")
    return tuple(bullets)


def _reference_bullets(references: Sequence[str], view: WorldView) -> tuple[str, ...]:
    """Templated citation lines (artifact title — id) from verified references."""
    bullets: list[str] = []
    for ref in references:
        node = view.get_node(ref)
        label = node.props.get("title", ref) if node is not None else ref
        bullets.append(f"{label} ({ref})")
    return tuple(bullets)


def _render_outline(slides: Sequence[Slide]) -> str:
    """Deterministic plain-text outline of the deck (what the tagger reasons over).

    One ``# Title`` line per slide, its subtitle (if any), then ``- bullet`` lines.
    This is the textual twin of the binary deck used for mention tagging + grounding.
    """
    lines: list[str] = []
    for slide in slides:
        lines.append(f"# {slide.title}")
        if slide.subtitle:
            lines.append(slide.subtitle)
        for bullet in slide.bullets:
            lines.append(f"- {bullet}")
        lines.append("")
    return "\n".join(lines).rstrip() + "\n"


# -- prompt assembly --------------------------------------------------------


def _system_prompt(kind: str) -> str:
    """Per-deck system prompt (cacheable across the run, §16.1)."""
    human = kind.replace("_", " ")
    return (
        f"You write the overview narrative for a {human} (a kickoff presentation) "
        "for an enterprise knowledge corpus. Write a short, factual overview "
        "grounded strictly in the supplied context. Refer only to the people and "
        "entities named in the roster, by exactly those names. Do not invent names, "
        "dates, or facts."
    )


def _stable_context(ctx: ProducerContext) -> list[str]:
    blocks: list[str] = []
    if ctx.company_profile:
        blocks.append(ctx.company_profile)
    if ctx.scenario_context:
        blocks.append(ctx.scenario_context)
    return blocks


def _task_brief(
    *,
    event: Event,
    kind: str,
    title: str,
    roster: Roster,
    authors: list[Node],
    reviewers: list[Node],
    candidates: Sequence[str],
) -> str:
    """The volatile per-deck suffix: brief + roster + templated facts (§16.1)."""
    topic = str(event.payload.get("topic") or event.payload.get("intent") or title)
    tone = event.payload.get("tone")
    lines = [
        f"Task: write the overview narrative for a {kind.replace('_', ' ')}.",
        f"Title: {title}.",
        f"Date: {event.timestamp.date().isoformat()}.",
        f"Topic: {topic}.",
    ]
    if tone:
        lines.append(f"Tone: {tone}.")
    if authors:
        lines.append(f"Prepared by: {_names(authors)}.")
    if reviewers:
        lines.append(f"Reviewed by: {_names(reviewers)}.")
    lines.append("")
    lines.append(roster.roster_block())
    if candidates:
        lines.append("")
        lines.append(
            "You may cite any of these prior artifacts by id, and you MUST report "
            "the ids you cite in references_used: " + ", ".join(candidates)
        )
    return "\n".join(lines)


def _repair_prompt(
    system: str,
    stable: list[str],
    brief: str,
    unresolved: Sequence[str],
) -> Prompt:
    """Re-assemble the prompt with a corrective instruction (the single repair, D30.3)."""
    correction = (
        "\n\nCORRECTION: your previous draft named "
        f"{', '.join(unresolved)}, which is not in scope. Rewrite using ONLY the "
        "roster names above; remove or replace every other name."
    )
    return assemble_prompt(system=system, stable_context=stable, brief=brief + correction)


# -- KG node/edge construction ----------------------------------------------


def _artifact_node(
    *,
    artifact_id: str,
    kind: str,
    medium: str,
    title: str,
    event: Event,
    authors: list[Node],
    reviewers: list[Node],
    path: str,
    slide_count: int,
) -> Node:
    """Build (or re-declare) the ``Artifact`` node this deck renders (§3)."""
    return Node(
        id=artifact_id,
        type=N_ARTIFACT,
        created_at=event.timestamp,
        props={
            "kind": kind,
            "medium": medium,
            "format": "pptx",
            "title": title,
            "path": path,
            "date": event.timestamp.date().isoformat(),
            "authors": [n.id for n in authors],
            "reviewers": [n.id for n in reviewers],
            "slide_count": slide_count,
            "event": event.id,
        },
        aliases=[title],
    )


def _build_edges(
    *,
    artifact_id: str,
    event: Event,
    view: WorldView,
    authors: list[Node],
    reviewers: list[Node],
    references: Sequence[str],
) -> list[Edge]:
    """Build the reified relationships the deck expresses (§3, §11.3, D16)."""
    edges: list[Edge] = []
    ts = event.timestamp
    for author in authors:
        edges.append(_edge(E_AUTHORED, author.id, artifact_id, ts))
    for reviewer in reviewers:
        edges.append(_edge(E_REVIEWED, reviewer.id, artifact_id, ts))
    for subject in event.subjects:
        if subject == artifact_id or view.get_node(subject) is None:
            continue
        edges.append(_edge(E_EXPRESSES, artifact_id, subject, ts))
    for ref in references:
        edges.append(_edge(E_REFERENCES, artifact_id, ref, ts))
    return edges


def _edge(edge_type: str, src: str, dst: str, ts: datetime) -> Edge:
    return Edge(
        id=f"edge:{edge_type}:{src}->{dst}",
        type=edge_type,
        src=src,
        dst=dst,
        created_at=ts,
    )


def _metadata(
    *,
    artifact_id: str,
    kind: str,
    medium: str,
    title: str,
    path: str,
    event: Event,
    authors: list[Node],
    reviewers: list[Node],
    references: Sequence[str],
    mentions: Sequence[Mention],
    issues: Sequence[ValidationIssue],
    slide_count: int,
) -> dict[str, object]:
    """The JSON twin of the deck (the ``pptx/json`` deliverable)."""
    return {
        "artifact_id": artifact_id,
        "kind": kind,
        "medium": medium,
        "format": "pptx",
        "title": title,
        "path": path,
        "event": event.id,
        "date": event.timestamp.date().isoformat(),
        "authors": [n.id for n in authors],
        "reviewers": [n.id for n in reviewers],
        "references": list(references),
        "slide_count": slide_count,
        "mentions": [m.to_dict() for m in mentions],
        "issues": [i.to_dict() for i in issues],
    }


# -- small helpers ----------------------------------------------------------


def _resolve_people(event: Event, view: WorldView, roles: Sequence[str]) -> list[Node]:
    """Resolve the in-scope person nodes bound to ``roles`` in ``event.actors``.

    Ids are de-duplicated preserving first-seen order; an id absent from the view is
    skipped (a producer can only name in-scope people).
    """
    seen: set[str] = set()
    people: list[Node] = []
    for role in roles:
        for pid in event.actors.get(role, []):
            if pid in seen:
                continue
            seen.add(pid)
            node = view.get_node(pid)
            if node is not None:
                people.append(node)
    return people


def _other_people(view: WorldView, *, exclude: set[str]) -> list[Node]:
    """In-scope people not already bound to a role, in (type, id) order (deterministic)."""
    return [n for n in view.nodes_by_type("Person") if n.id not in exclude]


def _candidate_references(view: WorldView, *, exclude: str) -> list[str]:
    """Recent in-scope artifact ids the model may cite (the candidate set, §16.3)."""
    return [n.id for n in view.nodes_by_type(N_ARTIFACT) if n.id != exclude]


def _names(nodes: Sequence[Node]) -> str:
    return ", ".join(_display_name(n) for n in nodes)


def _display_name(node: Node) -> str:
    name = node.props.get("name") or node.props.get("title")
    return str(name) if name else node.id.split(":", 1)[-1]


def _title_for(event: Event, kind: str) -> str:
    title = event.payload.get("title")
    if isinstance(title, str) and title:
        return title
    topic = event.payload.get("topic")
    human = kind.replace("_", " ").title()
    if isinstance(topic, str) and topic:
        return f"{human}: {topic}"
    return human


def _artifact_id_for(event: Event) -> str:
    """Use an existing abstract-artifact subject if present, else mint from the event."""
    for subject in event.subjects:
        if subject.startswith(("artifact:", "art:")):
            return subject
    return f"artifact:{event.id.split(':', 1)[-1]}"


def _slug(value: str) -> str:
    slug = re.sub(r"[^a-z0-9]+", "-", value.lower()).strip("-")
    return slug or "artifact"


# -- sample deck / CLI (human "opens cleanly in PowerPoint" check) ----------


def sample_deck() -> list[Slide]:
    """A canonical kickoff deck: title, agenda, overview, team, and references."""
    return [
        Slide(
            title="Payments Platform Kickoff",
            subtitle="Kickoff Deck · 2026-06-12 · Prepared by Ada Lovelace",
        ),
        Slide(title="Agenda", bullets=_AGENDA),
        Slide(
            title="Overview",
            bullets=(
                "The Payments Platform initiative kicks off this quarter.",
                "This deck frames the goals, the team, and the immediate next steps.",
            ),
        ),
        Slide(
            title="Team & Roles",
            bullets=("Ada Lovelace — Lead", "Alan Turing — Reviewer"),
        ),
        Slide(
            title="References",
            bullets=("Payments Design (artifact:design-pay)",),
        ),
    ]


def main(argv: Sequence[str] | None = None) -> int:
    """Write a sample kickoff ``.pptx`` so a human can confirm it opens in PowerPoint."""
    import sys

    args = list(sys.argv[1:] if argv is None else argv)
    out = args[0] if args else "kickoff_deck_sample.pptx"
    data = build_kickoff_deck(sample_deck())
    with open(out, "wb") as fh:
        fh.write(data)
    print(f"wrote {len(data)} bytes to {out}")
    return 0


# Self-register so the binding map can resolve "pptx" (registration fires on import,
# mirroring the playbook plugins). Guarded for idempotency: both the package
# ``__init__`` and the ``python -m enterprise_sim.producers.pptx`` entrypoint import
# this module, and a re-import must not raise a duplicate-registration error.
if "pptx" not in PRODUCERS:
    PRODUCERS.register(PptxProducer())


if __name__ == "__main__":
    raise SystemExit(main())
